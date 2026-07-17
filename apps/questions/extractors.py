# Formatos aceptados como material de origen. PDF permite además análisis nativo
# por Claude (Files API); PPTX/DOCX se procesan solo por su texto extraído.
ALLOWED_UPLOAD_EXTENSIONS = ('.pdf', '.pptx', '.docx')


def _extract_pdf_text(path) -> str:
    import pdfplumber

    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ''
            if text:
                parts.append(text)
    return '\n\n'.join(parts).strip()


def _extract_pptx_text(path) -> str:
    from pptx import Presentation

    parts = []
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(' | '.join(cells))
    return '\n'.join(parts).strip()


def _extract_docx_text(path) -> str:
    from docx import Document

    doc = Document(path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(' | '.join(cells))
    return '\n'.join(parts).strip()


def extract_document_text(path, filename) -> str:
    name = (filename or '').lower()
    if name.endswith('.pdf'):
        return _extract_pdf_text(path)
    if name.endswith('.pptx'):
        return _extract_pptx_text(path)
    if name.endswith('.docx'):
        return _extract_docx_text(path)
    return ''
